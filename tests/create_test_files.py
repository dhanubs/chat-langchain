import os
from pathlib import Path
from docx import Document
from pptx import Presentation
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter

# Create test files directory
test_files_dir = Path("tests/test_files")
test_files_dir.mkdir(parents=True, exist_ok=True)

# Create sample PDF
def create_sample_pdf():
    pdf_path = test_files_dir / "sample.pdf"
    c = canvas.Canvas(str(pdf_path), pagesize=letter)
    
    # Page 1
    c.drawString(100, 750, "Sample PDF Document")
    c.drawString(100, 700, "This is a test PDF file.")
    c.drawString(100, 650, "It contains multiple pages and text content.")
    
    # Page 2
    c.showPage()
    c.drawString(100, 750, "Page 2")
    c.drawString(100, 700, "This is the second page of the test document.")
    
    c.save()

# Create sample DOCX
def create_sample_docx():
    doc = Document()
    doc.add_heading('Sample DOCX Document', 0)
    
    # Add paragraphs
    doc.add_paragraph('This is a test Word document.')
    doc.add_paragraph('It contains text and a table.')
    
    # Add table
    table = doc.add_table(rows=3, cols=2)
    table.style = 'Table Grid'
    
    # Fill table
    cells = table.rows[0].cells
    cells[0].text = 'Header 1'
    cells[1].text = 'Header 2'
    
    cells = table.rows[1].cells
    cells[0].text = 'Data 1'
    cells[1].text = 'Data 2'
    
    cells = table.rows[2].cells
    cells[0].text = 'Data 3'
    cells[1].text = 'Data 4'
    
    doc.save(test_files_dir / "sample.docx")

# Create sample PPTX
def create_sample_pptx():
    prs = Presentation()
    
    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.shapes.placeholders[1]
    title.text = "Sample PowerPoint"
    subtitle.text = "Test Presentation"
    
    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    body = slide.shapes.placeholders[1]
    title.text = "Content Slide"
    tf = body.text_frame
    tf.text = "This is a test PowerPoint presentation"
    
    # Slide 3 with table
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shapes = slide.shapes
    
    rows, cols = 3, 2
    left = top = width = height = Inches(1.0)
    table = shapes.add_table(rows, cols, left, top, width, height).table
    
    table.rows[0].cells[0].text = 'Header 1'
    table.rows[0].cells[1].text = 'Header 2'
    table.rows[1].cells[0].text = 'Data 1'
    table.rows[1].cells[1].text = 'Data 2'
    table.rows[2].cells[0].text = 'Data 3'
    table.rows[2].cells[1].text = 'Data 4'
    
    prs.save(test_files_dir / "sample.pptx")

if __name__ == "__main__":
    print("Creating test files...")
    create_sample_pdf()
    create_sample_docx()
    create_sample_pptx()
    print(f"Test files created in {test_files_dir}") 