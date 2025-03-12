"""
Document processor module for handling various document types.
Supports PDF, DOCX, PPTX, and other document formats using direct document processing libraries:
- PyMuPDF (fitz) for PDF text extraction
- pdfplumber for table extraction from PDFs
- EasyOCR for image text extraction
- python-docx for Word documents
- python-pptx for PowerPoint presentations
"""

import os
import logging
import asyncio
import traceback
import warnings
from datetime import datetime
from pathlib import Path
from typing import List, Optional, Dict, Any, Union, Tuple, cast
import io
import numpy as np

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores.azuresearch import AzureSearch
from langchain_openai import AzureOpenAIEmbeddings

from backend.config import settings

import fitz  # PyMuPDF
from fitz import Document as FitzDocument
from fitz import Page as FitzPage
import pdfplumber
import easyocr
from PIL import Image
from docx import Document as DocxDocument
from pptx import Presentation

logger = logging.getLogger(__name__)

# Supported file extensions
SUPPORTED_EXTENSIONS = {
    ".pdf": "PDF Document",
    ".docx": "Word Document",
    ".doc": "Word Document (Legacy)",
    ".pptx": "PowerPoint Presentation",
    ".ppt": "PowerPoint Presentation (Legacy)",
    ".xlsx": "Excel Spreadsheet",
    ".xls": "Excel Spreadsheet (Legacy)",
    ".rtf": "Rich Text Format",
}

# Error categories for better error handling
class DocumentProcessingError(Exception):
    """Base exception for document processing errors."""
    def __init__(self, message: str, error_type: str, file_path: str, original_error: Optional[Exception] = None):
        self.message = message
        self.error_type = error_type
        self.file_path = file_path
        self.original_error = original_error
        super().__init__(self.message)

class UnsupportedFormatError(DocumentProcessingError):
    """Exception raised for unsupported file formats."""
    def __init__(self, file_path: str):
        super().__init__(
            f"Unsupported file format: {Path(file_path).suffix}",
            "unsupported_format",
            file_path
        )

class FileAccessError(DocumentProcessingError):
    """Exception raised for file access issues."""
    def __init__(self, file_path: str, original_error: Exception):
        super().__init__(
            f"Cannot access file: {original_error}",
            "file_access_error",
            file_path,
            original_error
        )

class ParsingError(DocumentProcessingError):
    """Exception raised for document parsing issues."""
    def __init__(self, file_path: str, original_error: Exception):
        super().__init__(
            f"Error parsing document: {original_error}",
            "parsing_error",
            file_path,
            original_error
        )

class EmptyContentError(DocumentProcessingError):
    """Exception raised when no content could be extracted."""
    def __init__(self, file_path: str):
        super().__init__(
            "No content could be extracted from the document",
            "empty_content",
            file_path
        )

class VectorStoreError(DocumentProcessingError):
    """Exception raised for vector store issues."""
    def __init__(self, file_path: str, original_error: Exception):
        super().__init__(
            f"Error adding document to vector store: {original_error}",
            "vector_store_error",
            file_path,
            original_error
        )

class DocumentProcessor:
    """
    Document processor class for handling various document types.
    Uses a combination of specialized libraries for each file type:
    - PDF: PyMuPDF for text, pdfplumber for tables, EasyOCR for images
    - DOCX: python-docx for text and tables
    - PPTX: python-pptx for slides and shapes
    
    Features:
    - Text extraction from PDFs, Word documents, and PowerPoint presentations
    - Table extraction from PDFs and Word documents
    - OCR for images in PDFs
    - Chunking of extracted text for better processing
    - Optional vector store integration with Azure Search
    - Parallel processing support for multiple files

    Note:
    On first initialization, EasyOCR will download its detection models (~45MB)
    to the ~/.EasyOCR/ directory. This is a one-time download per machine.
    """
    
    def __init__(
        self,
        chunk_size: int = 1000,
        chunk_overlap: int = 200,
        separators: List[str] = None,
        max_concurrency: int = 5,
        enable_ocr: bool = False,
        ocr_model_path: Optional[str] = None,
    ):
        """
        Initialize the document processor.
        
        Args:
            chunk_size: Size of text chunks for splitting
            chunk_overlap: Overlap between chunks
            separators: Custom separators for text splitting
            max_concurrency: Maximum number of files to process concurrently
            enable_ocr: Whether to enable OCR for images in PDFs
            ocr_model_path: Path to EasyOCR model directory (if None, uses default ~/.EasyOCR/)
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.separators = separators or ["\n\n", "\n", " ", ""]
        self.max_concurrency = max_concurrency
        self.enable_ocr = enable_ocr
        
        # Initialize text splitter
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=self.chunk_size,
            chunk_overlap=self.chunk_overlap,
            separators=self.separators,
            length_function=len
        )
        
        try:
            # Initialize vector store if Azure Search settings are available
            if settings.azure_search_service_name and settings.azure_search_key:
                embeddings = AzureOpenAIEmbeddings(
                    azure_deployment=settings.azure_openai_embedding_deployment,
                    openai_api_version=settings.azure_openai_api_version,
                    azure_endpoint=settings.azure_openai_endpoint,
                    api_key=settings.azure_openai_api_key,
                )
                
                self.vector_store = AzureSearch(
                    azure_search_endpoint=settings.azure_search_service_endpoint,
                    azure_search_key=settings.azure_search_admin_key,
                    index_name=settings.azure_search_index_name,
                    embedding_function=embeddings.embed_query,
                )
            else:
                self.vector_store = None
                logger.warning("Azure Search settings not configured. Vector store functionality will be disabled.")
            
            # Initialize EasyOCR only if enabled
            self.reader = None
            if enable_ocr:
                try:
                    logger.info("Initializing EasyOCR with%s model path", 
                              f" custom {ocr_model_path}" if ocr_model_path else " default")
                    self.reader = easyocr.Reader(['en'], gpu=False, 
                                               model_storage_directory=ocr_model_path)
                    logger.info("EasyOCR initialization complete")
                except Exception as ocr_error:
                    logger.warning(f"Failed to initialize EasyOCR: {str(ocr_error)}. OCR will be disabled.")
                    self.enable_ocr = False
            
        except Exception as e:
            logger.error("Failed to initialize DocumentProcessor.")
            logger.error(f"Error: {str(e)}")
            raise RuntimeError("DocumentProcessor initialization failed.") from e
    
    def is_supported_file(self, file_path: Union[str, Path]) -> bool:
        """
        Check if the file type is supported.
        
        Args:
            file_path: Path to the file
            
        Returns:
            bool: True if the file type is supported, False otherwise
        """
        file_path = Path(file_path)
        return file_path.suffix.lower() in SUPPORTED_EXTENSIONS
    
    def process_pdf(self, file_path: Union[str, Path]) -> str:
        """Extract text from PDF including OCR for images if enabled."""
        text_parts = []
        seen_images = set()  # Track image hashes to identify duplicates
        
        # Extract text and images with PyMuPDF
        with fitz.open(file_path) as doc:
            doc = cast(FitzDocument, doc)  # Type hint for better IDE support
            for page_num, page in enumerate(doc):
                page = cast(FitzPage, page)  # Type hint for better IDE support
                
                # Get text in reading order with proper formatting
                text = page.get_text("text")
                if text.strip():
                    text_parts.append(f"Page {page_num + 1}:\n{text}")
                
                # Handle images with OCR if enabled
                if self.enable_ocr and self.reader:
                    for img in page.get_images():
                        try:
                            xref = img[0]
                            base_img = doc.extract_image(xref)
                            
                            # Create a hash of the image data to identify duplicates
                            img_hash = hash(base_img["image"])
                            
                            # Skip if we've seen this image before (likely a logo or header/footer)
                            if img_hash in seen_images:
                                continue
                            
                            # Process new images
                            seen_images.add(img_hash)
                            image_bytes = base_img["image"]
                            image = Image.open(io.BytesIO(image_bytes))
                            
                            # Skip very small images (likely icons or decorative elements)
                            if image.width < 50 or image.height < 50:
                                continue
                                
                            image_np = np.array(image)
                            results = self.reader.readtext(image_np)
                            if results:
                                ocr_text = ' '.join([text for _, text, _ in results])
                                if ocr_text.strip():
                                    text_parts.append(f"Page {page_num + 1} (OCR):\n{ocr_text}")
                        except Exception as e:
                            logger.warning(f"Failed to process image on page {page_num + 1}: {str(e)}")
        
        return '\n\n'.join(text_parts)

    def process_docx(self, file_path: Union[str, Path]) -> str:
        """Extract text from DOCX files."""
        text_parts = []
        doc = DocxDocument(file_path)
        
        # Process paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Process tables
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = [cell.text.strip() for cell in row.cells]
                table_text.append('\t'.join(row_text))
            if table_text:
                text_parts.append('\n'.join(table_text))
        
        return '\n\n'.join(text_parts)

    def process_pptx(self, file_path: Union[str, Path]) -> str:
        """Extract text from PPTX files."""
        text_parts = []
        prs = Presentation(file_path)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
                
                # Handle tables
                if shape.has_table:
                    table_text = []
                    for row in shape.table.rows:
                        row_text = [cell.text.strip() for cell in row.cells]
                        table_text.append('\t'.join(row_text))
                    if table_text:
                        slide_text.append('\n'.join(table_text))
            
            if slide_text:
                text_parts.append(f"Slide {slide_num}:\n" + '\n'.join(slide_text))
        
        return '\n\n'.join(text_parts)

    def process_file(self, file_path: Union[str, Path]) -> Tuple[List[Document], Optional[DocumentProcessingError]]:
        """
        Process a single file and convert it to LangChain documents.
        
        Args:
            file_path: Path to the file
            
        Returns:
            Tuple[List[Document], Optional[DocumentProcessingError]]: 
                A tuple containing the list of documents and an optional error
        """
        file_path = Path(file_path)
        
        if not self.is_supported_file(file_path):
            error = UnsupportedFormatError(str(file_path))
            logger.warning(f"{error.error_type}: {error.message}")
            return [], error
        
        try:
            # Extract content based on file type
            if file_path.suffix.lower() in ['.pdf']:
                content = self.process_pdf(file_path)
            elif file_path.suffix.lower() in ['.docx', '.doc']:
                content = self.process_docx(file_path)
            elif file_path.suffix.lower() in ['.pptx', '.ppt']:
                content = self.process_pptx(file_path)
            else:
                return [], DocumentProcessingError(
                    f"Unsupported file format: {file_path.suffix}",
                    "unsupported_format",
                    str(file_path)
                )
            
            if not content or content.isspace():
                error = EmptyContentError(str(file_path))
                logger.warning(f"{error.error_type}: {error.message} for {file_path.name}")
                return [], error
            
            # Extract metadata
            metadata = {
                "source": str(file_path.name),
                "file_path": str(file_path),
                "file_type": SUPPORTED_EXTENSIONS[file_path.suffix.lower()],
                "title": file_path.stem,
                "created_at": datetime.utcnow().isoformat(),
                "chunk_size": self.chunk_size,
                "chunk_overlap": self.chunk_overlap,
                "processor_version": "3.0.0",
            }
            
            # Create a single document
            doc = Document(page_content=content, metadata=metadata)
            
            # Split the document
            documents = self.text_splitter.split_documents([doc])
            
            if not documents:
                error = EmptyContentError(str(file_path))
                logger.warning(f"{error.error_type}: No chunks created for {file_path.name}")
                return [], error
                
            return documents, None
            
        except Exception as e:
            error = self.classify_error(file_path, e)
            logger.error(f"{error.error_type} for {file_path.name}: {error.message}")
            logger.debug(f"Detailed error: {traceback.format_exc()}")
            return [], error
    
    def classify_error(self, file_path: Union[str, Path], exception: Exception) -> DocumentProcessingError:
        """
        Classify an exception into a specific DocumentProcessingError type.
        
        Args:
            file_path: Path to the file being processed
            exception: The exception that occurred
            
        Returns:
            DocumentProcessingError: A specific error type
        """
        file_path_str = str(file_path)
        
        # Check for file access errors
        if isinstance(exception, (FileNotFoundError, PermissionError, OSError)):
            return FileAccessError(file_path_str, exception)
        
        # Check for parsing errors
        if "parse" in str(exception).lower() or "convert" in str(exception).lower():
            return ParsingError(file_path_str, exception)
        
        # Default to a generic parsing error
        return ParsingError(file_path_str, exception)
    
    async def process_file_async(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Process a single file asynchronously and return processing results.
        
        Args:
            file_path: Path to the file
            
        Returns:
            Dict: Processing results including documents and statistics
        """
        try:
            # Process the file
            documents, error = self.process_file(file_path)
            
            # Handle processing errors
            if error:
                return {
                    "file_path": str(file_path),
                    "success": False,
                    "documents": [],
                    "error": error.message,
                    "error_type": error.error_type,
                    "file_name": Path(file_path).name
                }
            
            # Add to vector store
            try:
                if self.vector_store is None:
                    logger.warning(f"Vector store not configured. Skipping vector store addition for {Path(file_path).name}")
                else:
                    await self.vector_store.aadd_documents(documents)
            except Exception as e:
                error = VectorStoreError(str(file_path), e)
                logger.error(f"{error.error_type} for {Path(file_path).name}: {error.message}")
                return {
                    "file_path": str(file_path),
                    "success": False,
                    "documents": [],
                    "error": f"Failed to add to vector store: {str(e)}",
                    "error_type": "vector_store_error",
                    "file_name": Path(file_path).name
                }
            
            # Get file type
            file_type = SUPPORTED_EXTENSIONS[Path(file_path).suffix.lower()]
            
            logger.info(f"Processed {Path(file_path).name}: {len(documents)} chunks")
            
            return {
                "file_path": str(file_path),
                "success": True,
                "documents": documents,
                "chunks": len(documents),
                "file_type": file_type,
                "file_name": Path(file_path).name
            }
        
        except Exception as e:
            logger.error(f"Unexpected error processing {Path(file_path).name}: {str(e)}")
            logger.debug(f"Detailed error: {traceback.format_exc()}")
            return {
                "file_path": str(file_path),
                "success": False,
                "documents": [],
                "error": f"Unexpected error: {str(e)}",
                "error_type": "unexpected_error",
                "file_name": Path(file_path).name
            }
    
    async def process_directory(
        self,
        directory_path: Union[str, Path],
        recursive: bool = True,
        file_extensions: Optional[List[str]] = None,
        parallel: bool = True,
    ) -> Dict[str, Any]:
        """
        Process all supported files in a directory.
        
        Args:
            directory_path: Path to the directory
            recursive: Whether to search subdirectories
            file_extensions: List of file extensions to process (if None, process all supported types)
            parallel: Whether to process files in parallel
            
        Returns:
            Dict: Processing statistics
        """
        directory_path = Path(directory_path)
        
        # Validate directory path
        if not directory_path.exists():
            raise ValueError(f"Directory does not exist: {directory_path}")
        if not directory_path.is_dir():
            raise ValueError(f"Path is not a directory: {directory_path}")
        
        # Validate file extensions
        if file_extensions:
            invalid_extensions = [ext for ext in file_extensions if not ext.startswith(".")]
            if invalid_extensions:
                raise ValueError(f"Invalid file extensions (must start with '.'): {invalid_extensions}")
        
        # Filter extensions if provided
        extensions = file_extensions or list(SUPPORTED_EXTENSIONS.keys())
        
        # Find all matching files
        all_files = []
        try:
            if recursive:
                for ext in extensions:
                    all_files.extend(directory_path.rglob(f"*{ext}"))
            else:
                for ext in extensions:
                    all_files.extend(directory_path.glob(f"*{ext}"))
        except Exception as e:
            raise ValueError(f"Error searching for files: {str(e)}")
        
        # Check if any files were found
        if not all_files:
            return {
                "total_files": 0,
                "processed_files": 0,
                "failed_files": 0,
                "total_chunks": 0,
                "file_types": {},
                "parallel_processing": parallel,
                "max_concurrency": self.max_concurrency if parallel else 1,
                "warning": f"No matching files found in {directory_path}" + 
                          (f" with extensions {extensions}" if file_extensions else "")
            }
        
        # Process statistics
        stats = {
            "total_files": len(all_files),
            "processed_files": 0,
            "failed_files": 0,
            "total_chunks": 0,
            "file_types": {},
            "error_types": {},
            "parallel_processing": parallel,
            "max_concurrency": self.max_concurrency if parallel else 1,
            "start_time": datetime.utcnow().isoformat(),
            "failed_files_details": []
        }
        
        try:
            if not parallel:
                # Process files sequentially
                for file_path in all_files:
                    result = await self.process_file_async(file_path)
                    self._update_stats(stats, result)
            else:
                # Process files in parallel with controlled concurrency
                # Split files into batches to control memory usage
                batches = [all_files[i:i + self.max_concurrency] for i in range(0, len(all_files), self.max_concurrency)]
                
                for batch in batches:
                    # Process batch concurrently
                    results = await asyncio.gather(*[self.process_file_async(file_path) for file_path in batch])
                    
                    # Update statistics
                    for result in results:
                        self._update_stats(stats, result)
        except Exception as e:
            logger.error(f"Error during batch processing: {str(e)}")
            logger.debug(f"Detailed error: {traceback.format_exc()}")
            stats["batch_error"] = str(e)
        
        # Add end time
        stats["end_time"] = datetime.utcnow().isoformat()
        
        # Limit the number of failed file details to avoid huge responses
        if len(stats["failed_files_details"]) > 50:
            stats["failed_files_details"] = stats["failed_files_details"][:50]
            stats["failed_files_details_truncated"] = True
        
        return stats
    
    def _update_stats(self, stats: Dict[str, Any], result: Dict[str, Any]) -> None:
        """
        Update processing statistics with file processing result.
        
        Args:
            stats: Statistics dictionary to update
            result: File processing result
        """
        if result.get("success", False):
            stats["processed_files"] += 1
            stats["total_chunks"] += result.get("chunks", 0)
            
            # Update file type statistics
            file_type = result.get("file_type", "Unknown")
            if file_type not in stats["file_types"]:
                stats["file_types"][file_type] = 0
            stats["file_types"][file_type] += 1
        else:
            stats["failed_files"] += 1
            
            # Track error types
            error_type = result.get("error_type", "unknown_error")
            if error_type not in stats["error_types"]:
                stats["error_types"][error_type] = 0
            stats["error_types"][error_type] += 1
            
            # Add to failed files details
            stats["failed_files_details"].append({
                "file_name": result.get("file_name", "Unknown"),
                "file_path": result.get("file_path", "Unknown"),
                "error": result.get("error", "Unknown error"),
                "error_type": error_type
            }) 